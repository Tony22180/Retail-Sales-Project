from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import os
from datetime import datetime

def create_ppt_report():
    prs = Presentation()
    
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_frame = title.text_frame
    title_para = title_frame.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "RETAIL SALES PERFORMANCE"
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 32, 96)
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    subtitle_frame = subtitle.text_frame
    subtitle_para = subtitle_frame.add_paragraph()
    subtitle_run = subtitle_para.add_run()
    subtitle_run.text = f"Generated {datetime.now().strftime('%B %d, %Y')}"
    subtitle_run.font.size = Pt(18)
    subtitle_run.font.color.rgb = RGBColor(79, 129, 189)
    
    visualizations = [
        {
            "path": "reports/cohort_analysis.png",
            "title": "Customer Retention Analysis",
            "desc": "Monthly cohort retention rates showing customer loyalty trends"
        },
        {
            "path": "reports/sales_forecast.png",
            "title": "30-Day Sales Forecast",
            "desc": "ARIMA model projection of future sales performance"
        }
    ]
    
    for viz in visualizations:
        if not os.path.exists(viz["path"]):
            print(f"File not found: {viz['path']}")
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.5))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_run = title_para.add_run()
        title_run.text = viz["title"]
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        
        try:
            img = slide.shapes.add_picture(
                viz["path"],
                left=Inches((13.33-10)/2),
                top=Inches(1.5),
                width=Inches(10)
            )
        except Exception as e:
            print(f"Could not add image: {str(e)}")
            continue
            
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
        desc_frame = desc_box.text_frame
        desc_para = desc_frame.add_paragraph()
        desc_run = desc_para.add_run()
        desc_run.text = viz["desc"]
        desc_run.font.size = Pt(14)
        desc_run.font.color.rgb = RGBColor(100, 100, 100)
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "Key Customer Metrics"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    
    try:
        rfm_data = pd.read_excel("reports/rfm_analysis.xlsx")
        top_customers = rfm_data.nlargest(5, 'Monetary')
        
        table = slide.shapes.add_table(6, 3, Inches(1.5), Inches(1.5), Inches(10), Inches(3)).table
        
        table.cell(0, 0).text = "Customer ID"
        table.cell(0, 1).text = "Segment"
        table.cell(0, 2).text = "Total Spend"
        
        for i in range(5):
            table.cell(i+1, 0).text = str(top_customers.iloc[i]['Customer_ID'])
            table.cell(i+1, 1).text = top_customers.iloc[i]['Segment']
            table.cell(i+1, 2).text = f"${top_customers.iloc[i]['Monetary']:,.2f}"
        
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                if row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 32, 96)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True
                    
    except Exception as e:
        print(f"Could not create table: {str(e)}")
        error_box = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(6), Inches(1))
        error_para = error_box.text_frame.add_paragraph()
        error_run = error_para.add_run()
        error_run.text = "Data table not available"
    
    os.makedirs("presentation", exist_ok=True)
    output_path = "presentation/retail_insights.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint report generated: {output_path}")
    print(f"Total slides created: {len(prs.slides)}")

if __name__ == "__main__":
    create_ppt_report()